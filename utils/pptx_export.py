"""
amber/utils/pptx_export.py
───────────────────────────
Exportação para PowerPoint — idêntico à Célula 9 da versão Spyder.
Retorna bytes do .pptx para envio como HttpResponse de download.
Requer: pip install python-pptx
"""
from __future__ import annotations

import base64
import io

from .model import R_TARGET, P_TARGET


def export_pptx(result: dict, chart_b64: str) -> bytes:
    """
    Gera um slide PowerPoint com o estado atual da análise.

    Parâmetros
    ----------
    result    : dict retornado por model.full_analysis()
    chart_b64 : string base64 PNG do gráfico (retornado por charts.render_chart)

    Retorna
    -------
    bytes do arquivo .pptx
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("Instale python-pptx:  pip install python-pptx")

    def rgb(r, g, b): return RGBColor(r, g, b)

    a0     = result['alpha0'];  b0   = result['beta0']
    a_p    = result['alpha_p']; b_p  = result['beta_p']
    tm     = result['t_mission']
    beta_w = result['beta_w']
    n_obs  = result['n_obs']
    n_surv = result['n_surv'];  n_fail = result['n_fail']
    mu_po  = result['mu_post']; sig_po = result['sig_post']
    mu_pr  = result['mu_prior']
    p_pr   = result['p_prior']; p_po  = result['p_post']
    met    = result['met']
    sum_ks = result['sum_k_surv']
    sum_kf = result['sum_k_fail']

    # Decodifica a imagem base64 para bytes
    img_bytes = base64.b64decode(chart_b64)
    img_buf   = io.BytesIO(img_bytes)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(0xFA, 0xFA, 0xF8)

    # ── Cabeçalho ────────────────────────────────────────────────────────────
    hdr = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.58)
    )
    hdr.fill.solid(); hdr.fill.fore_color.rgb = rgb(0x18, 0x5F, 0xA5)
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.text = (
        f"AMBER v2.1  ·  T={tm:.1f}a  ·  Beta({a0:.1f},{b0:.1f})  ·  "
        f"{n_obs} obs. ({n_surv} cens., {n_fail} falhas)"
    )
    r = tf.paragraphs[0].runs[0]
    r.font.color.rgb = rgb(0xFF, 0xFF, 0xFF)
    r.font.size = Pt(12); r.font.bold = True
    tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.11)

    # ── Painel lateral ───────────────────────────────────────────────────────
    pan = slide.shapes.add_shape(
        1, Inches(0), Inches(0.58), Inches(3.10), Inches(6.92)
    )
    pan.fill.solid(); pan.fill.fore_color.rgb = rgb(0xF1, 0xEF, 0xE8)
    pan.line.color.rgb = rgb(0xD3, 0xD1, 0xC7); pan.line.width = Pt(0.5)

    # ── Imagem ───────────────────────────────────────────────────────────────
    slide.shapes.add_picture(img_buf, Inches(3.18), Inches(0.68),
                              Inches(10.0), Inches(6.55))

    # ── Texto do painel ──────────────────────────────────────────────────────
    s_clr = rgb(0x1D, 0x9E, 0x75) if met else rgb(0xD8, 0x5A, 0x30)
    rows = [
        ("Prior",               True,  rgb(0x2C,0x2C,0x2A), 11),
        ("",                    False, rgb(0x88,0x87,0x80),   7),
        (f"α₀ = {a0:.1f}",     True,  rgb(0x18,0x5F,0xA5),  16),
        (f"β₀ = {b0:.1f}",     True,  rgb(0x18,0x5F,0xA5),  16),
        (f"E[R]|prior = {mu_pr:.4f}",  False, rgb(0x5F,0x5E,0x5A), 10),
        (f"P[R≥R*]|prior = {p_pr:.3f}", False, rgb(0x5F,0x5E,0x5A), 10),
        ("",                    False, rgb(0x88,0x87,0x80),   7),
        ("Modelo",              True,  rgb(0x2C,0x2C,0x2A),  11),
        ("",                    False, rgb(0x88,0x87,0x80),   5),
        (f"T_MISSION = {tm:.1f} anos", True, rgb(0x18,0x5F,0xA5), 13),
        (f"β_W = {beta_w:.1f}", False, rgb(0x5F,0x5E,0x5A),  10),
        ("",                    False, rgb(0x88,0x87,0x80),   7),
        ("Observações",         True,  rgb(0x2C,0x2C,0x2A),  11),
        ("",                    False, rgb(0x88,0x87,0x80),   5),
        (f"Total = {n_obs}",    False, rgb(0x5F,0x5E,0x5A),  10),
        (f"Cens. = {n_surv}  (Σk={sum_ks:.3f})", False, rgb(0x5F,0x5E,0x5A), 10),
        (f"Falhas = {n_fail}  (Σk={sum_kf:.3f})",
         False, rgb(0xD8,0x5A,0x30) if n_fail else rgb(0x5F,0x5E,0x5A), 10),
        ("",                    False, rgb(0x88,0x87,0x80),   7),
        ("Posterior",           True,  rgb(0x2C,0x2C,0x2A),  11),
        ("",                    False, rgb(0x88,0x87,0x80),   5),
        (f"α = {a_p:.3f}   β = {b_p:.3f}", False, rgb(0x5F,0x5E,0x5A), 10),
        (f"E[R] = {mu_po:.4f}   σ = {sig_po:.4f}", False, rgb(0x5F,0x5E,0x5A), 10),
        ("",                    False, rgb(0x88,0x87,0x80),   7),
        (f"P[R≥{R_TARGET:.2f}] = {p_po:.3f}", True, rgb(0x2C,0x2C,0x2A), 14),
        ("✓  Alvo atingido" if met else "✗  Abaixo do alvo", True, s_clr, 11),
    ]
    tb_shape = slide.shapes.add_textbox(
        Inches(0.14), Inches(0.72), Inches(2.90), Inches(6.65)
    )
    tf2 = tb_shape.text_frame; tf2.word_wrap = True
    for i, (txt, bold, clr, sz) in enumerate(rows):
        pp = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        pp.text = txt; pp.space_after = Pt(1)
        if pp.runs:
            rr = pp.runs[0]
            rr.font.bold = bold; rr.font.color.rgb = clr; rr.font.size = Pt(sz)

    # ── Rodapé ───────────────────────────────────────────────────────────────
    ft = slide.shapes.add_textbox(
        Inches(0), Inches(7.27), Inches(13.33), Inches(0.23)
    )
    ft.text_frame.text = (
        f"AMBER v2.1  ·  Ceerma–UFPE  ·  "
        f"Prior Beta({a0},{b0})  ·  T={tm:.1f}a  ·  "
        f"Alvo: P[R≥{R_TARGET}]≥{P_TARGET}"
    )
    ft.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    ft.text_frame.paragraphs[0].runs[0].font.color.rgb = rgb(0x88, 0x87, 0x80)
    ft.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
